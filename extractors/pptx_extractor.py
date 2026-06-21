from pptx import Presentation

def extract_text(path):
    prs = Presentation(path)
    text_runs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_runs.append(f"--- Slide {slide_num} ---")
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text_runs.append(run.text)
                            
    return "\n".join(text_runs)